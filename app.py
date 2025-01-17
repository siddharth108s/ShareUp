import os
import bcrypt
import json
import mysql.connector
from flask import Flask, render_template, request, redirect, url_for, session, flash, send_file
from Crypto.Cipher import AES
from Crypto.Util.Padding import pad, unpad
from Crypto.Hash import SHA256
from Crypto.PublicKey import DSA
from Crypto.Signature import DSS
from flask import abort
import logging
from datetime import datetime, timedelta
import re
from functools import wraps
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
import io

app = Flask(__name__)
app.secret_key = 'your_secret_key'
app.config['UPLOAD_FOLDER'] = 'uploads'  # Make sure this directory exists


db_config = {
'host' : "127.0.0.1",
'user': "root",
'password' : "",
'database' : "is_project",
'port' : 3306
}


def init_db():
    conn = mysql.connector.connect(**db_config)
    return conn


def close_db(conn):
    conn.close()


def read_aes_key(filename):
    with open(filename, 'rb') as file:
        key = file.read(16)  # AES-128 requires a 16-byte key
    return key


def aes_encrypt(key, file_data):
    cipher = AES.new(key, AES.MODE_CBC)
    iv = cipher.iv
    ciphertext = cipher.encrypt(pad(file_data, AES.block_size))
    return iv, ciphertext


def aes_decrypt(key, iv, ciphertext):
    cipher = AES.new(key, AES.MODE_CBC, iv)
    decrypted = unpad(cipher.decrypt(ciphertext), AES.block_size)
    return decrypted


def hash_keyword(keyword):
    hasher = SHA256.new()
    hasher.update(keyword.encode('utf-8'))
    return hasher.hexdigest()

def sign(private_key_path, data):
    with open(private_key_path, 'rb') as key_file:
        private_key = DSA.import_key(key_file.read())
    data_hash = SHA256.new(data)
    signer = DSS.new(private_key, 'fips-186-3')
    signature = signer.sign(data_hash)
    return signature


def verify(public_key_path, data, signature):
    with open(public_key_path, 'rb') as key_file:
        public_key = DSA.import_key(key_file.read())
    data_hash = SHA256.new(data)
    verifier = DSS.new(public_key, 'fips-186-3')
    try:
        verifier.verify(data_hash, signature)
        return True
    except (ValueError, TypeError):
        return False


def index_file(file_id, keywords):
    conn = init_db()
    cursor = conn.cursor()
    for keyword in keywords:
        hashed_keyword = hash_keyword(keyword)
        try:
            insert_query = "INSERT INTO FileIndex (keyword, file_id) VALUES (%s, %s)"
            cursor.execute(insert_query, (hashed_keyword, file_id))
        except mysql.connector.Error as err:
            print(f"Error inserting keyword '{keyword}': {err}")
    conn.commit()
    cursor.close()
    close_db(conn)


def upload_file(file_path, aes_key_path, private_key_path):
    key = read_aes_key(aes_key_path)
    with open(file_path, 'rb') as file:
        file_data = file.read()
    iv, ciphertext = aes_encrypt(key, file_data)
    signature = sign(private_key_path, ciphertext)
    log_action("File Signature", f"File: {os.path.basename(file_path)}, Signature: {signature.hex()}")

    # Get the file size
    filesize = os.path.getsize(file_path)  # Get the size of the file in bytes

    conn = init_db()
    cursor = conn.cursor()

    target_username = request.form['username']
    sender_username = session['username']
    filename = os.path.basename(file_path)

    cursor.execute("SELECT username FROM Users WHERE username = %s", (target_username,))
    result = cursor.fetchone()

    if not result:
        cursor.close()
        close_db(conn)
        return 0

    insert_query = "INSERT INTO Files (username, sender, filename, iv, encrypted_file, signature, filesize, upload_time) VALUES (%s, %s, %s, %s, %s, %s, %s, NOW() + INTERVAL 0 HOUR + INTERVAL 0 MINUTE)"
    try:
        cursor.execute(insert_query, (target_username, sender_username, filename, iv, ciphertext, signature, filesize))
        file_id = cursor.lastrowid
        conn.commit()

        # List of text-based file extensions
        text_extensions = ('.txt', '.py', '.c', '.cpp', '.h', '.java', '.js', '.html', 
                          '.css', '.md', '.json', '.xml', '.csv', '.pdf', '.docx', '.pptx')
        
        # Index content for text-based files
        if filename.lower().endswith(text_extensions):
            try:
                content = ""
                file_ext = os.path.splitext(filename)[1].lower()
                
                if file_ext == '.pdf':
                    # Extract text from PDF
                    with open(file_path, 'rb') as pdf_file:
                        pdf_reader = PdfReader(pdf_file)
                        for page in pdf_reader.pages:
                            content += page.extract_text() + "\n"
                
                elif file_ext == '.docx':
                    # Extract text from DOCX
                    doc = Document(file_path)
                    for para in doc.paragraphs:
                        content += para.text + "\n"
                
                elif file_ext == '.pptx':
                    # Extract text from PPTX
                    prs = Presentation(file_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                
                else:
                    # Original text file handling
                    encodings = ['utf-8', 'latin-1', 'ascii']
                    for encoding in encodings:
                        try:
                            with open(file_path, 'r', encoding=encoding) as f:
                                content = f.read()
                            break
                        except UnicodeDecodeError:
                            continue
                
                if content:
                    # Split content into words, filtering out special characters and empty strings
                    keywords = set(word.strip() for word in re.split(r'[\s\n\t\r.,;:!?(){}[\]<>"`\']+', content) if word.strip())
                    index_file(file_id, keywords)

            except Exception as e:
                log_action("Indexing Error", f"Failed to index file {filename}: {str(e)}")

        if cursor.lastrowid:
            log_action("File Upload", f"File ID: {file_id}, From: {sender_username}, To: {target_username}, Filename: {filename}, Filesize: {filesize} bytes")
        return 1
    except mysql.connector.Error:
        return 0
    finally:
        cursor.close()
        close_db(conn)


def receive_files(username):
    log_action("Files Retrieved", f"User {username} accessed their received files")
    conn = init_db()
    cursor = conn.cursor()
    select_query = "SELECT file_id, sender, filename, iv, encrypted_file, signature, upload_time,filesize FROM Files WHERE username = %s"
    cursor.execute(select_query, (username,))
    results = cursor.fetchall()

    verified_files = []
    for file in results:
        file_id, sender, filename, iv, encrypted_file, signature, upload_time, filesize = file
        verified = verify("dsa_public_key.pem", encrypted_file, signature)
        adjusted_time = upload_time - timedelta(hours=0, minutes=0)
        verified_files.append((file_id, sender, filename, iv, encrypted_file, signature, verified, adjusted_time, filesize))

    cursor.close()
    close_db(conn)
    return verified_files


def search_encrypted_files(keyword):
    log_action("Search", f"Admin {session.get('username')} searched for keyword: {keyword}")
    conn = init_db()
    cursor = conn.cursor()

    hashed_keyword = hash_keyword(keyword)

    cursor.execute("SELECT file_id FROM FileIndex WHERE keyword = %s", (hashed_keyword,))
    file_ids = cursor.fetchall()

    cursor.close()
    close_db(conn)

    return [file_id[0] for file_id in file_ids]


def check_credentials(username, password):
    conn = init_db()
    cursor = conn.cursor()
    cursor.execute("SELECT password, role FROM Users WHERE username = %s", (username,))
    result = cursor.fetchone()
    close_db(conn)

    if result:
        hashed_password, role = result
        if bcrypt.checkpw(password.encode('utf-8'), hashed_password.encode('utf-8')):
            return role
    return None


@app.route('/')
def home():
    return render_template('login.html')


@app.route('/login', methods=['POST'])
def login():
    username = request.form['username']
    password = request.form['password']

    role = check_credentials(username, password)

    if role:
        log_action("Login", f"User '{username}' logged in successfully with role '{role}'")
        session['username'] = username
        session['role'] = role
        if role == 'admin':
            return redirect(url_for('admin_page'))
        else:
            return redirect(url_for('user_page'))
    else:
        log_action("Failed Login", f"Failed login attempt for username '{username}'")
        flash("Invalid credentials, please try again.", "error")  # Flash the error message
        return redirect(url_for('home'))


@app.route('/signup', methods=['GET', 'POST'])
def signup():
    if request.method == 'POST':
        username = request.form['username']
        password1 = request.form['password1']
        password2 = request.form['password2']

        if password1 != password2:
            flash("Passwords do not match", "error")
            return redirect(url_for('signup'))

        hashed_password = bcrypt.hashpw(password1.encode('utf-8'), bcrypt.gensalt())

        conn = init_db()
        cursor = conn.cursor()
        insert_query = """
            INSERT INTO Users (username, password, role)
            VALUES (%s, %s, 'user')
        """
        try:
            cursor.execute(insert_query, (username, hashed_password.decode('utf-8')))
            conn.commit()
            log_action("Signup", f"New user account created: {username}")
            flash("Account created successfully. Please log in.", "success")
            return redirect(url_for('home'))
        except mysql.connector.Error as err:
            log_action("Failed Signup", f"Failed to create account for {username}: {str(err)}")
            flash("Error creating account: " + str(err), "error")
            return redirect(url_for('signup'))
        finally:
            cursor.close()
            close_db(conn)

    return render_template('signup.html')


@app.route('/user', methods=['GET', 'POST'])
def user_page():
    message = None
    if request.method == 'POST':
        if 'file' not in request.files:
            message = {"type": "error", "text": "No file selected"}
            return render_template('user.html', message=message)
            
        file = request.files['file']
        if file.filename == '':
            message = {"type": "error", "text": "No file selected"}
            return render_template('user.html', message=message)

        try:
            # Remove the logic for checking file size and creating the uploads directory
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            with open(file_path, 'wb') as f:
                f.write(file.read())

            result = upload_file(file_path, "aes_key.bin", "dsa_private_key.pem")

            # Clean up the temporary file
            if os.path.exists(file_path):
                os.remove(file_path)

            if result == 1:
                message = {"type": "success", "text": "File sent successfully!"}
            else:
                message = {"type": "error", "text": "Failed to send file. Invalid username or file limit exceeded."}

        except Exception as e:
            log_action("Upload Error", f"Error uploading file: {str(e)}")
            message = {"type": "error", "text": f"Error uploading file: {str(e)}"}

    # Get received files
    username = session.get('username')
    files = receive_files(username)
    sent_files = get_sent_files(username)
    
    return render_template('user.html', message=message, files=files, sent_files=sent_files)


@app.route('/received_files')
def received_files():
    username = session.get('username')
    files = receive_files(username)
    return render_template('received_files.html', files=files)


@app.route('/download_file/<int:file_id>')
def download_file(file_id):
    conn = init_db()
    cursor = conn.cursor()
    cursor.execute("SELECT iv, encrypted_file, filename FROM Files WHERE file_id = %s", (file_id,))
    result = cursor.fetchone()
    close_db(conn)

    if result:
        iv, encrypted_file, filename = result
        key = read_aes_key("aes_key.bin")
        decrypted_data = aes_decrypt(key, iv, encrypted_file)

        temp_path = f'temp_decrypted_file_{file_id}{os.path.splitext(filename)[1]}'
        with open(temp_path, 'wb') as f:
            f.write(decrypted_data)

        log_action("File Download", f"File ID: {file_id}, Filename: {filename}, Downloaded by: {session.get('username')}")
        return send_file(temp_path, as_attachment=True, download_name=filename)

    return "File not found."


def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not session.get('role') == 'admin':
            abort(403)  # Forbidden access
        return f(*args, **kwargs)
    return decorated_function


@app.route('/admin', methods=['GET', 'POST'])
@admin_required
def admin_page():
    conn = init_db()
    cursor = conn.cursor()
    
    # Fetch data from all three tables
    cursor.execute("SELECT * FROM Users")
    users = cursor.fetchall()
    
    cursor.execute("SELECT file_id, filename, username, sender, upload_time, filesize FROM Files")
    files = cursor.fetchall()
    
    cursor.execute("SELECT * FROM FileIndex")
    file_index = cursor.fetchall()
    
    # Calculate total storage
    cursor.execute("SELECT SUM(filesize) FROM Files")
    total_storage = cursor.fetchone()[0] or 0
    
    cursor.close()
    close_db(conn)
    
    if request.method == 'POST':
        keyword = request.form['keyword']
        file_ids = search_encrypted_files(keyword)
        return render_template('admin.html', keyword=keyword, file_ids=file_ids, users=users, files=files, file_index=file_index, total_storage=total_storage)
    
    return render_template('admin.html', users=users, files=files, file_index=file_index, total_storage=total_storage)


@app.route('/delete_user', methods=['POST'])
@admin_required
def delete_user():
    username = request.form['username']
    conn = init_db()
    cursor = conn.cursor()

    try:
        # First check if user exists
        cursor.execute("SELECT username FROM Users WHERE username = %s", (username,))
        if not cursor.fetchone():
            flash(f"User '{username}' does not exist.", "error")
            return redirect(url_for('admin_page'))

        cursor.execute("SELECT file_id FROM Files WHERE username = %s", (username,))
        file_ids = [row[0] for row in cursor.fetchall()]

        if file_ids:
            cursor.executemany("DELETE FROM FileIndex WHERE file_id = %s", [(file_id,) for file_id in file_ids])

        cursor.execute("DELETE FROM Files WHERE username = %s", (username,))

        cursor.execute("DELETE FROM Users WHERE username = %s", (username,))

        conn.commit()
        log_action("User Deletion", f"User '{username}' and associated files deleted by admin {session.get('username')}")
        flash(f"User '{username}' and associated files have been deleted.", "success")
    except mysql.connector.Error as err:
        log_action("Failed User Deletion", f"Failed to delete user '{username}': {str(err)}")
        flash(f"Error deleting user: {err}", "error")
    finally:
        cursor.close()
        close_db(conn)

    return redirect(url_for('admin_page'))


@app.route('/logout')
def logout():
    username = session.get('username')
    log_action("Logout", f"User '{username}' logged out")
    session.pop('username', None)
    session.pop('role', None)
    return redirect(url_for('home'))


def log_action(action, details):
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    log_entry = f"[{timestamp}] {action}: {details}\n"
    with open('logs.txt', 'a') as log_file:
        log_file.write(log_entry)


def get_sent_files(username):
    log_action("Files Retrieved", f"User {username} accessed their sent files")
    conn = init_db()
    cursor = conn.cursor()
    select_query = "SELECT file_id, username, filename, iv, encrypted_file, signature, upload_time, filesize FROM Files WHERE sender = %s"
    cursor.execute(select_query, (username,))
    results = cursor.fetchall()

    verified_files = []
    for file in results:
        file_id, recipient, filename, iv, encrypted_file, signature, upload_time, filesize = file
        verified = verify("dsa_public_key.pem", encrypted_file, signature)
        adjusted_time = upload_time - timedelta(hours=0, minutes=0)
        verified_files.append((file_id, recipient, filename, iv, encrypted_file, signature, verified, adjusted_time, filesize))

    cursor.close()
    close_db(conn)
    return verified_files


@app.route('/delete_file', methods=['POST'])
@admin_required
def delete_file():
    file_id = request.form['file_id']
    conn = init_db()
    cursor = conn.cursor()

    try:
        # First delete from FileIndex table
        cursor.execute("DELETE FROM FileIndex WHERE file_id = %s", (file_id,))
        
        # Then delete from Files table
        cursor.execute("DELETE FROM Files WHERE file_id = %s", (file_id,))
        
        conn.commit()
        log_action("File Deletion", f"File ID {file_id} deleted by admin {session.get('username')}")
        flash(f"File ID {file_id} has been deleted.", "success")
    except mysql.connector.Error as err:
        log_action("Failed File Deletion", f"Failed to delete file ID {file_id}: {str(err)}")
        flash(f"Error deleting file: {err}", "error")
    finally:
        cursor.close()
        close_db(conn)

    return redirect(url_for('admin_page'))


@app.errorhandler(413)
def request_entity_too_large(error):
    flash('File too large. Maximum size is 16MB.', 'error')
    return redirect(url_for('user_page')), 413


if __name__ == "__main__":
    app.run(debug=True)
